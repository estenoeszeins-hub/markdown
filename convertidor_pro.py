import customtkinter as ctk
from tkinter import filedialog, messagebox, simpledialog
import mammoth
import pandas as pd
from pptx import Presentation
from markdownify import markdownify as md
import os

# --- CONFIGURACIÓN DE ACCESO ---
CONTRASEÑA_MAESTRA = "Chris_PAss2026MKD@"
# Ruta donde se guarda la activación para que no pida pass cada vez
PATH_ACTIVA = os.path.join(os.environ.get("APPDATA", ""), "activacion_uarm_mkd.dat")

def verificar_acceso():
    """Pide la contraseña solo la primera vez."""
    if os.path.exists(PATH_ACTIVA):
        return True

    # Ventana temporal para la contraseña
    temp_root = ctk.CTk()
    temp_root.withdraw()
    
    password_input = simpledialog.askstring(
        "Acceso Restringido", 
        "Introduce la contraseña de Chris para usar la herramienta:",
        show='*', # Oculta los caracteres
        parent=None
    )
    
    if password_input == CONTRASEÑA_MAESTRA:
        try:
            with open(PATH_ACTIVA, "w") as f:
                f.write("ACCESO_PERMITIDO")
            messagebox.showinfo("Éxito", "Acceso concedido. El programa se ha activado en este equipo.")
            temp_root.destroy()
            return True
        except:
            return False
    else:
        messagebox.showerror("Error", "Contraseña incorrecta. Contacta con Christopher Ccoicca.")
        temp_root.destroy()
        return False

# --- CLASE PRINCIPAL DE LA APP ---
ctk.set_appearance_mode("System")
ctk.set_default_color_theme("blue")

class App(ctk.CTk):
    def __init__(self):
        super().__init__()

        self.title("Convertidor MD - Christopher Ccoicca")
        self.geometry("700x630")

        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(3, weight=1)

        self.label = ctk.CTkLabel(self, text="Compresor de Documentos a Markdown", font=("Roboto", 22, "bold"))
        self.label.grid(row=0, column=0, pady=(20, 10), sticky="ew")

        self.btn_select = ctk.CTkButton(self, text="Seleccionar Archivos", 
                                       command=self.select_files, height=40, font=("Roboto", 14))
        self.btn_select.grid(row=1, column=0, pady=10)

        self.path_label = ctk.CTkLabel(self, text="Listo para procesar", text_color="gray")
        self.path_label.grid(row=2, column=0, pady=5)

        self.textbox = ctk.CTkTextbox(self, width=600, height=300, font=("Consolas", 11))
        self.textbox.grid(row=3, column=0, pady=15, padx=20, sticky="nsew")

        self.btn_save = ctk.CTkButton(self, text="Convertir y Guardar Todo", 
                                     state="disabled", command=self.process_and_save_all, 
                                     fg_color="#2ecc71", hover_color="#27ae60")
        self.btn_save.grid(row=4, column=0, pady=(10, 5))

        self.footer_label = ctk.CTkLabel(self, text="Created by: Christopher Ccoicca", 
                                         font=("Roboto", 10, "italic"), text_color="gray")
        self.footer_label.grid(row=5, column=0, pady=(5, 10), padx=20, sticky="se")

        self.selected_paths = []

    def select_files(self):
        self.selected_paths = filedialog.askopenfilenames(
            filetypes=[("Documentos de Oficina", "*.docx *.xlsx *.pptx")]
        )
        if self.selected_paths:
            self.path_label.configure(text=f"{len(self.selected_paths)} archivos seleccionados", text_color="#3498db")
            self.textbox.delete("1.0", "end")
            for path in self.selected_paths:
                self.textbox.insert("end", f"• {os.path.basename(path)}\n")
            self.btn_save.configure(state="normal")

    def process_file(self, path):
        ext = path.split(".")[-1].lower()
        try:
            if ext == "docx":
                with open(path, "rb") as f:
                    return md(mammoth.convert_to_html(f).value)
            elif ext == "xlsx":
                df = pd.read_excel(path)
                return md(df.to_html(index=False))
            elif ext == "pptx":
                prs = Presentation(path)
                return "\n\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        except Exception as e:
            return f"Error en {os.path.basename(path)}: {str(e)}"
        return ""

    def process_and_save_all(self):
        if not self.selected_paths: return
        target_dir = filedialog.askdirectory(title="Carpeta de destino")
        if target_dir:
            success_count = 0
            for path in self.selected_paths:
                content = self.process_file(path)
                base_name = os.path.splitext(os.path.basename(path))[0]
                save_path = os.path.join(target_dir, f"{base_name}.md")
                try:
                    with open(save_path, "w", encoding="utf-8") as f:
                        f.write(content)
                    success_count += 1
                except Exception as e:
                    print(f"Error: {e}")

            messagebox.showinfo("Completado", f"Se han convertido {success_count} archivos.")
            self.textbox.delete("1.0", "end")
            self.btn_save.configure(state="disabled")

if __name__ == "__main__":
    import ctypes
    try: ctypes.windll.shcore.SetProcessDpiAwareness(1)
    except: pass
        
    # Verificar contraseña antes de abrir la App
    if verificar_acceso():
        app = App()
        app.mainloop()
