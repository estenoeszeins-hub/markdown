import streamlit as st
import mammoth
import pandas as pd
from pptx import Presentation
from markdownify import markdownify as md
import base64
from io import BytesIO
from PIL import Image
import os

# --- CONFIGURACIÓN ---
st.set_page_config(page_title="Zeins Native Pro", page_icon="⚡", layout="wide")

CONTRASEÑA_MAESTRA = "Chris_PAss2026MKD@"

if 'autenticado' not in st.session_state:
    st.session_state['autenticado'] = False

# --- LOGIN ---
if not st.session_state['autenticado']:
    st.title("🔐 Acceso Restringido")
    st.markdown("Plataforma de extracción semántica")
    pass_input = st.text_input("Contraseña:", type="password")
    if st.button("Entrar") and pass_input == CONTRASEÑA_MAESTRA:
        st.session_state['autenticado'] = True
        st.rerun()
    elif pass_input:
        st.error("Credenciales incorrectas.")
    st.stop()

# --- FUNCIONES DE EXTRACCIÓN ---
def extract_image(shape):
    """Extrae, limpia y codifica imágenes ignorando transparencias problemáticas."""
    try:
        img = Image.open(BytesIO(shape.image.blob))
        # Limpieza de modo color
        if img.mode in ("RGBA", "P"): 
            img = img.convert("RGB")
        # Redimensionar para optimizar peso (máximo 800px)
        img.thumbnail((800, 800))
        buf = BytesIO()
        img.save(buf, format="JPEG", quality=80)
        encoded = base64.b64encode(buf.getvalue()).decode()
        return f"\n\n![Elemento Visual](data:image/jpeg;base64,{encoded})\n\n"
    except Exception as e:
        return f"\n> [!AVISO] Error al renderizar imagen: {e}\n"

def procesar_presentacion(file, extraer_imagenes):
    """Procesa el PPTX usando Bounding Box Mapping para mantener el orden lógico."""
    prs = Presentation(file)
    md_output = [f"# Documento: {file.name}\n"]
    
    for idx, slide in enumerate(prs.slides):
        md_output.append(f"\n---\n## 🗂️ DIAPOSITIVA {idx+1}\n")
        
        # 1. MAPEADO ESPACIAL
        elementos_visuales = []
        for shape in slide.shapes:
            if not hasattr(shape, "top") or not hasattr(shape, "left"):
                continue
            elementos_visuales.append({
                "top": shape.top,
                "left": shape.left,
                "shape": shape
            })
            
        # Orden de lectura: Arriba -> Abajo, Izquierda -> Derecha
        elementos_visuales.sort(key=lambda x: (x["top"], x["left"]))
        
        # 2. EXTRACCIÓN ORDENADA
        for item in elementos_visuales:
            shape = item["shape"]
            
            # Procesar Tablas
            if shape.has_table:
                md_output.append("\n")
                for row_idx, row in enumerate(shape.table.rows):
                    fila_texto = [cell.text_frame.text.replace('\n',' ').replace('\r','') for cell in row.cells]
                    md_output.append("| " + " | ".join(fila_texto) + " |")
                    # Cabecera de la tabla MD
                    if row_idx == 0:
                        md_output.append("|" + "|".join(["---"] * len(row.cells)) + "|")
                md_output.append("\n")
                
            # Procesar Texto y Títulos
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    texto = paragraph.text.strip()
                    if texto:
                        # Detectar niveles de viñetas
                        if paragraph.level > 0:
                            md_output.append(f"{'  ' * paragraph.level}* {texto}")
                        else:
                            # Detectar texto en negrita como posibles subtítulos
                            if paragraph.runs and len(paragraph.runs) > 0 and paragraph.runs[0].font.bold:
                                md_output.append(f"### {texto}")
                            else:
                                md_output.append(f"{texto}\n")
                                
            # Procesar Imágenes (Si el usuario lo permite)
            elif extraer_imagenes and shape.shape_type == 13: # 13 = PICTURE
                md_output.append(extract_image(shape))
                
        # 3. NOTAS DEL EXPOSITOR
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            md_output.append(f"\n> 💡 **Notas del orador:** {slide.notes_slide.notes_text_frame.text.strip()}\n")
            
    return "\n".join(md_output)

# --- INTERFAZ PRINCIPAL ---
st.title("⚙️ Motor de Extracción Semántica")
st.markdown("Transforma documentos a formato nativo para IA mediante mapeo de coordenadas espaciales.")

extraer_imagenes = st.checkbox("🖼️ Extraer elementos visuales (Aumenta el tamaño del archivo)", value=True)

archivos = st.file_uploader("Arrastra tus archivos aquí", type=["pptx", "docx", "xlsx"], accept_multiple_files=True)

if archivos:
    for f in archivos:
        try:
            with st.spinner(f"Mapeando coordenadas de {f.name}..."):
                if f.name.endswith(".pptx"):
                    resultado = procesar_presentacion(f, extraer_imagenes)
                elif f.name.endswith(".docx"):
                    resultado = md(mammoth.convert_to_html(f).value)
                elif f.name.endswith(".xlsx"):
                    resultado = md(pd.read_excel(f).to_html(index=False))
                    
                with st.expander(f"✅ {f.name} procesado"):
                    st.download_button(f"📥 Descargar Markdown ({f.name})", resultado, file_name=f"{f.name}.md")
        except Exception as e:
            st.error(f"Error procesando {f.name}: {e}")

st.sidebar.markdown("---")
st.sidebar.success("Sistema Activo - Zeins Edition")
